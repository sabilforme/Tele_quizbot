import os
from typing import Optional

PDF_OK = True
DOCX_OK = True
PPTX_OK = True

try:
    from pdfminer.high_level import extract_text as pdf_extract_text
except Exception:
    PDF_OK = False
try:
    import docx
except Exception:
    DOCX_OK = False
try:
    from pptx import Presentation
except Exception:
    PPTX_OK = False

from ocr import ocr_space_file


async def extract_text_any(path: str, suffix: str, lang: str) -> str:
    suffix = (suffix or "").lower()

    if suffix == ".pdf":
        # 1) جرّب النص الأصلي
        if PDF_OK:
            try:
                text = pdf_extract_text(path) or ""
                if len(text.strip()) > 300:
                    return text
            except Exception:
                pass
        # 2) OCR عبر OCR.space (يدعم PDF متعدد الصفحات)
        try:
            ocr_text = await ocr_space_file(path, lang)
            if ocr_text:
                return ocr_text
        except Exception:
            pass
        return ""

    if suffix == ".docx" and DOCX_OK:
        try:
            d = docx.Document(path)
            return "\n".join(p.text for p in d.paragraphs)
        except Exception:
            return ""

    if suffix == ".pptx" and PPTX_OK:
        try:
            prs = Presentation(path)
            chunks = []
            for slide in prs.slides:
                # نص الشرائح
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        chunks.append("\n".join([p.text for p in shape.text_frame.paragraphs]))
                # ملاحظات المحاضر إن وجدت
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    chunks.append(slide.notes_slide.notes_text_frame.text)
            return "\n".join(chunks)
        except Exception:
            return ""

    if suffix == ".txt":
        try:
            return open(path, "r", encoding="utf-8", errors="ignore").read()
        except Exception:
            return ""

    # صور: jpg/png/tiff … → OCR
    try:
        return await ocr_space_file(path, lang)
    except Exception:
        return ""
